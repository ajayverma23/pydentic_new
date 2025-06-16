
import pdfkit
import markdown2
import base64
import requests
import re
import csv
from datetime import datetime
import os
from pptx import Presentation
from pptx.util import Inches, Pt

from app_schema.data_validation import ContentRequest



# markdown_to_html
def markdown_to_html(md_text):
    return markdown2.markdown(md_text)

# html_to_pdf
def html_to_pdf(html_text, output_pdf_path):
    path_wkhtmltopdf = r"C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe"
    config = pdfkit.configuration(wkhtmltopdf=path_wkhtmltopdf)
    options = {"enable-local-file-access": True}
    pdfkit.from_string(html_text, output_pdf_path, configuration=config, options=options)

# img_url_to_base64
def img_url_to_base64(img_url):
    response = requests.get(img_url)
    img_format = img_url.split('.')[-1]
    base64_str = base64.b64encode(response.content).decode('utf-8')
    return f"data:image/{img_format};base64,{base64_str}"

def img_file_to_base64(image_path):
    with open(image_path, "rb") as img_file:
        base64_str = base64.b64encode(img_file.read()).decode('utf-8')
    img_format = image_path.split('.')[-1]
    return f"data:image/{img_format};base64,{base64_str}"


def extract_url_from_output(text):
    if not isinstance(text, str):
        return None
    match = re.search(r'(https?://\S+\.(?:jpg|jpeg|png|gif))', text)
    return match.group(1) if match else None

# extract short subject or title from your content to save pdf file
def slugify(text, max_length=30):
    text = re.sub(r'[^\w\s-]', '', text).strip().lower()
    text = re.sub(r'[-\s]+', '_', text)
    return text[:max_length]


# extract only unique items / remove duplicate tag / ref / hashtag etc
def dedupe_lines(text):
    seen = set()
    lines = []
    for line in text.splitlines():
        if line.strip() and line not in seen:
            seen.add(line)
            lines.append(line)
    return "\n".join(lines)

def log_status(log_path, level, agent, task, status, output=None, error=None):
    # If agent is an object with .role, use it; otherwise, use the string directly
    agent_role = getattr(agent, "role", agent)
    task_desc = getattr(task, "description", task)
    file_exists = os.path.isfile(log_path)
    with open(log_path, 'a', newline='') as logfile:
        writer = csv.writer(logfile)
        if not file_exists:
            writer.writerow(["Timestamp", "Level", "Agent", "Task", "Status", "Output Length", "Error"])
        writer.writerow([
            datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            level,
            agent_role,
            task_desc,
            status,
            len(output) if output else None,
            error
        ])

def generate_ppt_from_content_request(content_request, pptx_path):
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = content_request.topic or "Presentation"
    slide.placeholders[1].text = f"By: {content_request.role}"

    # Introduction Slide
    intro_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(intro_slide_layout)
    slide.shapes.title.text = "Introduction"
    slide.placeholders[1].text = (
        f"{content_request.context}\n\n"
        f"Target Audience: {content_request.target_audience} ({content_request.audience_level})\n"
        f"Content Type: {content_request.content_type}\n"
        f"Subject: {content_request.subject}"
    )

    # Main Content Slide
    main_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(main_slide_layout)
    slide.shapes.title.text = "Content"
    slide.placeholders[1].text = (
        f"{content_request.user_input}\n\n"
        f"Reference: {content_request.reference_data}"
    )

    # Save the presentation
    prs.save(pptx_path)
    print(f"Presentation saved to {pptx_path}")


import csv
from datetime import datetime
import os

def log_inputs(original: ContentRequest, improved: dict, log_path):
    logs_dir = r"D:\Aj\GenAI\pydentic_new\output\logs"
    os.makedirs(logs_dir, exist_ok=True)
    
    #log_path = os.path.join(logs_dir, "input_review_logs.csv")
    
    # Create/append to CSV with headers if missing
    header = [
        "timestamp", 
        "original_user_input", 
        "improved_user_input",
        "original_context", 
        "improved_context",
        "original_subject", 
        "improved_subject",
        "original_topic", 
        "improved_topic",
        "original_reference_data", 
        "improved_reference_data",
        "content_type",
        "target_audience"
    ]
    
    row = {
        "timestamp": datetime.now().isoformat(),
        "original_user_input": original.user_input,
        "improved_user_input": improved.get("user_input", ""),
        "original_context": original.context,
        "improved_context": improved.get("context", ""),
        "original_subject": original.subject,
        "improved_subject": improved.get("subject", ""),
        "original_topic": original.topic,
        "improved_topic": improved.get("topic", ""),
        "original_reference_data": original.reference_data,
        "improved_reference_data": improved.get("reference_data", ""),
        "content_type": original.content_type,
        "target_audience": original.target_audience
    }
    
    file_exists = os.path.isfile(log_path)
    with open(log_path, "a", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=header)
        if not file_exists:
            writer.writeheader()
        writer.writerow(row)

