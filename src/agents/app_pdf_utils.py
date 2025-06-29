
import pdfkit
import markdown2
import base64
import requests
import re
import csv
from datetime import datetime
import os
from pptx import Presentation
from pptx.util import Inches, Pt


#set PYTHONPATH="D:\Aj\GenAI\pydentic_new\src;%PYTHONPATH%"

from app_schema.data_validation import ContentRequest

#from data_validation import ContentRequest



# markdown_to_html
def markdown_to_html(md_text):
    return markdown2.markdown(md_text)

# html_to_pdf
def html_to_pdf(html_text, output_pdf_path):
    path_wkhtmltopdf = r"C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe"
    config = pdfkit.configuration(wkhtmltopdf=path_wkhtmltopdf)
    options = {"enable-local-file-access": True}
    pdfkit.from_string(html_text, output_pdf_path, configuration=config, options=options)

# img_url_to_base64
def img_url_to_base64(img_url):
    response = requests.get(img_url)
    img_format = img_url.split('.')[-1]
    base64_str = base64.b64encode(response.content).decode('utf-8')
    return f"data:image/{img_format};base64,{base64_str}"

def img_file_to_base64(image_path):
    with open(image_path, "rb") as img_file:
        base64_str = base64.b64encode(img_file.read()).decode('utf-8')
    img_format = image_path.split('.')[-1]
    return f"data:image/{img_format};base64,{base64_str}"


def extract_url_from_output(text):
    if not isinstance(text, str):
        return None
    match = re.search(r'(https?://\S+\.(?:jpg|jpeg|png|gif))', text)
    return match.group(1) if match else None

# extract short subject or title from your content to save pdf file
def slugify(text, max_length=30):
    text = re.sub(r'[^\w\s-]', '', text).strip().lower()
    text = re.sub(r'[-\s]+', '_', text)
    return text[:max_length]


# extract only unique items / remove duplicate tag / ref / hashtag etc
def dedupe_lines(text):
    seen = set()
    lines = []
    for line in text.splitlines():
        if line.strip() and line not in seen:
            seen.add(line)
            lines.append(line)
    return "\n".join(lines)

# def log_status(log_path, level, agent, task, status, output=None, error=None):
#     # If agent is an object with .role, use it; otherwise, use the string directly
    
#     agent_role = getattr(agent, "role", agent)
#     task_desc = getattr(task, "description", task)
#     file_exists = os.path.isfile(log_path)
#     with open(log_path, 'a', newline='') as logfile:
#         writer = csv.writer(logfile)
#         if not file_exists:
#             writer.writerow(["Timestamp", "Level", "Agent", "Task", "Status", "Output Length", "Error"])
#         writer.writerow([
#             datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
#             level,
#             agent_role,
#             task_desc,
#             status,
#             len(output) if output else None,
#             error
#         ])


def log_status(log_path, level, agent, task, status, output=None, error=None):
    """Logs task status to a CSV file using utf-8 encoding."""
    agent_role = agent.role if hasattr(agent, 'role') else 'Unknown Agent'
    task_description = task.description if hasattr(task, 'description') else 'Unknown Task'
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    log_data = [timestamp, 
                level, 
                agent_role, 
                task_description, 
                status,
                len(output) if output else None,
                error]

    file_exists = os.path.isfile(log_path)
    try:
        with open(log_path, mode='a', newline='', encoding='utf-8') as logfile:
            writer = csv.writer(logfile)
            if not file_exists:
                #writer.writerow(["Timestamp", "Level", "Agent", "Task", "Status"])
                writer.writerow(["Timestamp", "Level", "Agent", "Task", "Status", "Output Length", "Error"])
            writer.writerow(log_data)
    except Exception as e:
        print(f"Error writing to log file: {e}")

def generate_ppt_from_content_request(content_request, pptx_path):
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = content_request.topic or "Presentation"
    slide.placeholders[1].text = f"By: {content_request.role}"

    # Introduction Slide
    intro_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(intro_slide_layout)
    slide.shapes.title.text = "Introduction"
    slide.placeholders[1].text = (
        f"{content_request.context}\n\n"
        f"Target Audience: {content_request.target_audience} ({content_request.audience_level})\n"
        f"Content Type: {content_request.content_type}\n"
        f"Subject: {content_request.subject}"
    )

    # Main Content Slide
    main_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(main_slide_layout)
    slide.shapes.title.text = "Content"
    slide.placeholders[1].text = (
        f"{content_request.user_input}\n\n"
        f"Reference: {content_request.reference_data}"
    )

    # Save the presentation
    prs.save(pptx_path)
    print(f"Presentation saved to {pptx_path}")


import csv
from datetime import datetime
import os

def log_inputs(original: ContentRequest, improved: dict, log_path):
    logs_dir = r"D:\Aj\GenAI\pydentic_new\output\logs"
    os.makedirs(logs_dir, exist_ok=True)
    
    #log_path = os.path.join(logs_dir, "input_review_logs.csv")
    
    # Create/append to CSV with headers if missing
    header = [
        "timestamp", 
        "original_user_input", 
        "improved_user_input",
        "original_context", 
        "improved_context",
        "original_subject", 
        "improved_subject",
        "original_topic", 
        "improved_topic",
        "original_reference_data", 
        "improved_reference_data",
        "content_type",
        "target_audience"
    ]
    
    row = {
        "timestamp": datetime.now().isoformat(),
        "original_user_input": original.user_input,
        "improved_user_input": improved.get("user_input", ""),
        "original_context": original.context,
        "improved_context": improved.get("context", ""),
        "original_subject": original.subject,
        "improved_subject": improved.get("subject", ""),
        "original_topic": original.topic,
        "improved_topic": improved.get("topic", ""),
        "original_reference_data": original.reference_data,
        "improved_reference_data": improved.get("reference_data", ""),
        "content_type": original.content_type,
        "target_audience": original.target_audience
    }
    
    file_exists = os.path.isfile(log_path)
    with open(log_path, "a", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=header)
        if not file_exists:
            writer.writeheader()
        writer.writerow(row)


HARDCODED_PROMPTS = {
    "Blog": {
        "system_prompt": "You are a professional blog writer. Write engaging, informative blogs tailored to the specified audience.",
        "user_prompt": "Create a blog post that is clear, concise, and relevant to the topic provided."
    },
    "Training Contents": {
        "system_prompt": "You are an expert training content creator. Develop comprehensive and structured training material.",
        "user_prompt": "Generate training content with examples, exercises, and clear explanations."


    },
    # ...other content types
}


#not used yet
def generate_system_prompt(content_type: str, domain: str = None, user_system_prompt: str = None):
    """
    Generates a dynamic system prompt, combining hardcoded prompts with user-provided prompts (if available).

    Args:
        content_type: The type of content to generate (e.g., "Blog", "Training Contents").
        domain: The specific domain for the content (optional, e.g., "healthcare", "cloud").
        user_system_prompt: A custom system prompt provided by the user (optional).

    Returns:
        A string representing the combined system prompt.
    """
    hardcoded_prompt = ""
    if content_type == "Training Contents":
        default_domain = "General AI" if domain is None else domain
        hardcoded_prompt = f"""You are an expert AI educational content generator specializing in {default_domain} education, with deep knowledge of both {default_domain.lower()} and artificial intelligence. Your task is to create engaging, progressive training content for professionals learning about AI and its applications in {default_domain.lower()}. Follow content generation guidelines to build all aspects of education. Follow all other instruction correctly and remember to follow all directions to the T and don't be afraid to get creative"""
    elif content_type == "Blog":
        hardcoded_prompt = "You are a world-class blogger specializing in cutting-edge technologies. You write compelling, informative, and engaging blog posts for a technical audience."
    elif content_type == "Test cases":
        hardcoded_prompt = "You are a world-class software QA specializing in creating various test cases."
    elif content_type == "Requirement":
        hardcoded_prompt = "You are a world-class software architect specializing in creating business requirements."
    elif content_type == "Process":
        hardcoded_prompt = "You are a world-class business analyst specializing in creating business processes."
    else:
        hardcoded_prompt = "You are a helpful AI assistant."

    if user_system_prompt:
        combined_prompt = f"{user_system_prompt}\n{hardcoded_prompt}"
        return combined_prompt
    else:
        return hardcoded_prompt
    
#not used yet
def generate_user_prompt(content_type: str,topic : str=None,user_user_prompt: str = None):
        hardcoded_prompt = ""
        if content_type == "Training Contents":
            hardcoded_prompt = f"""You have expert knowladge in GenAI frameworks and provide excellent training"""
        elif content_type == "Blog":
            hardcoded_prompt = "You can create effective SEO blog and can get top rating"
        elif content_type == "Test cases":
            hardcoded_prompt = "You can write test cases with all different cases"
        elif content_type == "Requirement":
            hardcoded_prompt = "Can create requiremnt documents with perfect details."
        elif content_type == "Process":
            hardcoded_prompt = "Able to create process document very well"
        else:
            hardcoded_prompt = "You are a helpful AI assistant."
        if user_user_prompt:
            combined_prompt = f"{user_user_prompt}\n{hardcoded_prompt}"
            return combined_prompt
        else:
            return hardcoded_prompt



def get_final_prompts(content_request, HARDCODED_PROMPTS):
    ct = content_request.content_type
    hardcoded = HARDCODED_PROMPTS.get(ct, {"system_prompt": "", "user_prompt": ""})

    # Merge user-provided prompts if present
    system_prompt = (hardcoded["system_prompt"] + " " + content_request.system_prompt.strip()) if content_request.system_prompt else hardcoded["system_prompt"]
    user_prompt = (hardcoded["user_prompt"] + " " + content_request.user_prompt.strip()) if content_request.user_prompt else hardcoded["user_prompt"]

    return system_prompt.strip(), user_prompt.strip()

def serialize_result(result):
    # Try common serialization methods
    if hasattr(result, "model_dump"):
        return result.model_dump()
    elif hasattr(result, "dict"):  # For Pydantic v1
        return result.dict()
    elif isinstance(result, str):
        return result
    else:
        return str(result)


